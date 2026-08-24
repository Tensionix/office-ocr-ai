#!/usr/bin/env python3
from __future__ import annotations

import argparse
import gc
import hashlib
import os
import re
import shutil
import sys
import tempfile
import threading
import time
from pathlib import Path
from typing import List, Tuple

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from tqdm import tqdm

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from project_paths import (
    BASE_DIR,
    CONFIG_DIR,
    CONFIG_API_KEY_GEMINI_FILE,
    CONFIG_API_KEY_OPENAI_FILE,
    API_KEY_GEMINI_FILE,
    API_KEY_OPENAI_FILE,
    LEGACY_CONFIG_API_KEY_GEMINI_FILE,
    LEGACY_CONFIG_API_KEY_OPENAI_FILE,
    LLM_SETTINGS_FILE,
    TEMP_DIR,
    ensure_project_dirs,
    iter_project_files,
    normalized_source_relative_path,
    output_path_for,
)
from providers.gemini_provider import call_markdown_vision as gemini_call_markdown_vision
from providers.openai_provider import call_markdown_vision as openai_call_markdown_vision

try:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tif", ".tiff", ".webp"}
LOCAL_EXTS = {".docx", ".pptx", ".xlsx", ".html", ".txt", ".csv", ".pdf"}
AI_EXTS = {".pdf"} | IMAGE_EXTS
CAMERA_EXIF_TAGS = {271, 272, 274, 306, 34853, 36867, 36868}
DEFAULT_GEMINI_MODEL_CHAIN = ["gemini-3.5-flash", "gemini-2.5-pro", "gemini-3.1-pro-preview"]

DEFAULT_OCR_PROMPT = """
Ты — профессиональный OCR-ассистент и технический редактор.
Нужно преобразовать документ в аккуратный Markdown.
Правила:
1. Сохраняй структуру заголовков и разделов.
2. Таблицы переносить как корректные Markdown-таблицы.
3. Списки сохранять списками.
4. Не выдумывать содержимое; только то, что реально видно.
5. Игнорировать служебный мусор вроде повторяющихся колонтитулов и номеров страниц, если они не несут смысловой нагрузки.
6. Если в батче несколько страниц, верни единый непрерывный Markdown-фрагмент в их порядке.
7. По возможности правильно угадывай абзацы, вложенные списки, нумерацию, кодовые блоки, цитаты и таблицы.
8. Верни только Markdown без пояснений.
9. Если это фото листа A4 с перспективным искажением, мысленно выровняй документ и восстанавливай текст по листу, а не по наклонённой фотографии.
10. Основной язык документа может быть русским; также возможны английские и немецкие фрагменты.
11. Не заменяй кириллицу похожими латинскими буквами и не латинизируй русские слова.
12. Для русских ФИО, должностей, названий организаций и географических названий особенно внимательно проверяй среднюю и конечную часть слов.
13. Обращай особое внимание на типичные OCR-путаницы: ц/щ, з/э, е/ё, и/й, ь/ъ, а также на смешение кириллицы и латиницы.
14. Не добавляй символы, которых нет на изображении, и не делай строку визуально длиннее, чем она выглядит в документе.
15. Если документ или фрагмент состоит почти полностью из числовых данных, считай приоритетом точность цифр и разделителей, а не угадывание буквенных слов.
16. Если слово частично закрыто, смазано или потеряно примерно не более чем на 25%, и по видимой части плюс контексту его можно восстановить однозначно, допускается аккуратно дорезолвить это слово.
17. Если однозначности нет, ничего не додумывай и оставляй только то, что реально можно уверенно прочитать.
18. Именованные сущности критически важны: фамилии, имена, отчества, названия организаций, городов, населённых пунктов, областей, стран, рек, гор и других объектов не нормализуй в более привычную форму без явного визуального основания.
19. Если редкая фамилия или топоним видны достаточно уверенно, сохраняй именно редкое написание, даже если более распространённый вариант кажется модели знакомее.
20. Не склоняй автоматически украинские фамилии, немецкие фамилии и фамилии, оканчивающиеся на согласную букву, если такое склонение не подтверждается самим изображением.
21. Для таких фамилий предпочитай исходную форму, даже если в контексте ожидается косвенный падеж.
""".strip()

DEFAULT_OCR_REFINE_PROMPT = """
Ты выполняешь второй проход OCR-коррекции по тем же изображениям документа.
Ниже уже есть черновик Markdown после первого прохода.

Задача:
1. Сохранить структуру Markdown, порядок блоков, списков и таблиц.
2. Исправить только ошибки распознавания, сверяясь с изображениями.
3. Для русскоязычного текста считай, что инициалы и первая заглавная буква фамилии или названия часто уже распознаны верно; особенно перепроверяй среднюю и конечную часть слов.
4. Уделяй особое внимание ошибкам: ц/щ, з/э, е/ё, и/й, ь/ъ, а также латинским буквам внутри кириллических слов.
5. Не сокращай текст, не пересказывай, не стилизуй заново и не меняй смысл.
6. Если в документе встречаются английские или немецкие фрагменты, сохраняй их как есть.
7. Не добавляй новые символы и не удлиняй строку по сравнению с тем, что реально видно на изображении.
8. Если документ numeric-heavy, не подмешивай латиницу в числовые токены и не заменяй цифры похожими буквами вроде 1/l или 5/S.
9. Если слово частично закрыто, смазано или потеряно примерно не более чем на 25%, и по изображению плюс контексту его можно восстановить однозначно, допускается аккуратно исправить это место.
10. Если уверенности нет, не додумывай символы и не расширяй слово сверх того, что можно надёжно вывести из изображения и контекста.
11. Именованные сущности критически важны: фамилии, имена, отчества, названия организаций, городов, населённых пунктов, областей, стран, рек, гор и других объектов не нормализуй в более привычную форму без явного визуального основания.
12. Если редкая фамилия, топоним или название выглядят правдоподобно и подтверждаются изображением, предпочитай исходное редкое написание более популярной догадке.
13. Не склоняй автоматически украинские фамилии, немецкие фамилии и фамилии, оканчивающиеся на согласную букву, если такое склонение не подтверждается изображением.
14. Для таких фамилий предпочитай исходную форму, даже если в контексте ожидается косвенный падеж.
15. Верни только исправленный Markdown без пояснений.

Языки документа: {language_hint}
Документ: {document_name}
Батч страниц: {chunk_index} из {chunk_count}
Режим numeric-heavy: {numeric_mode_hint}
Подсказки по именованным сущностям:
{entity_hints}

Черновик первого прохода:
```md
{draft_markdown}
```
""".strip()


_GEMINI_CLIENTS = threading.local()


def read_text_file(path: Path) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp1251", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except Exception:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def write_text_file(path: Path, text: str) -> None:
    ensure_parent_dir(path)
    path.write_text(text.rstrip() + "\n", encoding="utf-8")


def resolve_config_file(setting_value: str, default_name: str) -> Path:
    raw = str(setting_value or "").strip() or default_name
    path = Path(raw)
    if path.is_absolute():
        return path
    return CONFIG_DIR / path


def ensure_prompt_template(path: Path, default_text: str) -> str:
    if not path.exists():
        write_text_file(path, default_text)
    return read_text_file(path).strip()


def render_prompt_template(template: str, values: dict[str, str]) -> str:
    rendered = template
    for key, value in values.items():
        rendered = rendered.replace("{" + key + "}", value)
    return rendered.strip()


def contains_cyrillic(text: str) -> bool:
    return bool(re.search(r"[А-Яа-яЁё]", text or ""))


def cyrillic_ratio(text: str) -> float:
    letters = re.findall(r"[A-Za-zА-Яа-яЁё]", text or "")
    if not letters:
        return 0.0
    cyrillic = re.findall(r"[А-Яа-яЁё]", text)
    return len(cyrillic) / len(letters)


def numeric_signal_ratio(text: str) -> float:
    chars = [ch for ch in (text or "") if not ch.isspace()]
    if not chars:
        return 0.0
    numeric_chars = 0
    for ch in chars:
        if ch.isdigit() or ch in ".,;:+-–—()[]{}\\/|=_%#№":
            numeric_chars += 1
    return numeric_chars / len(chars)


def _parse_scalar_value(value: str):
    text = value.strip()
    if not text:
        return ""
    lowered = text.lower()
    if lowered == "true":
        return True
    if lowered == "false":
        return False
    if (text.startswith('"') and text.endswith('"')) or (text.startswith("'") and text.endswith("'")):
        return text[1:-1]
    return text


def parse_entity_hints_yaml(text: str) -> dict[str, list]:
    result: dict[str, list] = {}
    current_section: str | None = None
    current_item: dict | None = None
    current_list_key: str | None = None

    for raw_line in text.splitlines():
        if not raw_line.strip():
            continue
        stripped = raw_line.lstrip()
        if stripped.startswith("#"):
            continue
        indent = len(raw_line) - len(stripped)

        if indent == 0 and stripped.endswith(":"):
            current_section = stripped[:-1].strip()
            result.setdefault(current_section, [])
            current_item = None
            current_list_key = None
            continue

        if current_section is None:
            continue

        section_items = result[current_section]

        if indent == 2 and stripped.startswith("- "):
            current_list_key = None
            item_text = stripped[2:].strip()
            if ":" in item_text:
                key, value = item_text.split(":", 1)
                current_item = {key.strip(): _parse_scalar_value(value)}
                section_items.append(current_item)
            else:
                current_item = None
                section_items.append(_parse_scalar_value(item_text))
            continue

        if indent == 4 and current_item is not None and stripped.endswith(":"):
            key = stripped[:-1].strip()
            current_item[key] = []
            current_list_key = key
            continue

        if indent == 4 and current_item is not None and ":" in stripped:
            key, value = stripped.split(":", 1)
            current_item[key.strip()] = _parse_scalar_value(value)
            current_list_key = None
            continue

        if indent == 6 and current_item is not None and current_list_key and stripped.startswith("- "):
            current_item[current_list_key].append(_parse_scalar_value(stripped[2:]))
            continue

    return result


def format_entity_hints_text(data: dict[str, list]) -> str:
    lines: list[str] = []
    for section, entries in data.items():
        if not entries:
            continue
        lines.append(f"[{section}]")
        for entry in entries:
            if isinstance(entry, str):
                lines.append(f"- {entry}")
                continue
            if isinstance(entry, dict):
                lemma = str(entry.get("lemma", "")).strip()
                indeclinable = entry.get("indeclinable")
                forms = entry.get("forms", [])
                bits: list[str] = []
                if lemma:
                    bits.append(f"lemma={lemma}")
                if indeclinable is True:
                    bits.append("indeclinable=true")
                if isinstance(forms, list) and forms:
                    bits.append("forms=" + ", ".join(str(item).strip() for item in forms if str(item).strip()))
                if bits:
                    lines.append("- " + "; ".join(bits))
                    continue
            lines.append(f"- {str(entry).strip()}")
        lines.append("")
    return "\n".join(lines).strip()


DEFAULT_ENTITY_HINTS = """
families:
  - lemma: Марушкина
    forms:
      - Марушкина
      - Марушкиной
      - Марушкину
      - Марушкине
  - lemma: Мезер
    indeclinable: true
    forms:
      - Мезер

places:
  - Омск
  - Нефтеюганск
  - Благовещенск

regions:
  - ХМАО - Югра
  - Нефтеюганский район

organizations:
  - ООО "ИТП "ГРАД""

water_and_mountains:
  - Иртыш
""".strip()


def md_escape_cell(value: str) -> str:
    return value.replace("\r", " ").replace("\n", " ").strip().replace("|", r"\|")


def ensure_parent_dir(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def docx_to_md(docx_path: Path) -> str:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(docx_path))

    def iter_block_items(parent):
        parent_elm = parent.element.body
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    out_lines: List[str] = []
    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            text = (block.text or "").strip()
            if not text:
                continue
            style = (block.style.name if block.style is not None else "") or ""
            style_lower = style.lower()
            if style_lower.startswith("heading"):
                level = 1
                parts = style.split()
                if len(parts) >= 2 and parts[-1].isdigit():
                    level = max(1, min(6, int(parts[-1])))
                out_lines.extend(("#" * level + " " + text, ""))
                continue
            is_list = False
            try:
                ppr = block._p.pPr
                if ppr is not None and ppr.numPr is not None:
                    is_list = True
            except Exception:
                pass
            out_lines.append(f"- {text}" if is_list else text)
            out_lines.append("")
            continue

        rows = []
        for row in block.rows:
            rows.append([md_escape_cell(cell.text or "") for cell in row.cells])
        if not rows:
            continue
        col_count = max(1, len(rows[0]))
        header = rows[0] + [""] * (col_count - len(rows[0]))
        out_lines.append("| " + " | ".join(header[:col_count]) + " |")
        out_lines.append("| " + " | ".join(["---"] * col_count) + " |")
        for row in rows[1:]:
            padded = row + [""] * (col_count - len(row))
            out_lines.append("| " + " | ".join(padded[:col_count]) + " |")
        out_lines.append("")
    return "\n".join(out_lines).strip() + "\n"


def pptx_to_md(pptx_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    out_lines: List[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        out_lines.extend((f"## Slide {index}", ""))
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text or "").strip()
            if not text:
                continue
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            if len(lines) == 1:
                out_lines.append(lines[0])
            else:
                out_lines.extend(f"- {line}" for line in lines)
            out_lines.append("")
    return "\n".join(out_lines).strip() + "\n"


def xlsx_to_md(xlsx_path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(xlsx_path), data_only=True)
    out_lines: List[str] = []
    for sheet in wb.worksheets:
        out_lines.extend((f"## Sheet: {sheet.title}", ""))
        used_rows: List[List[str]] = []
        for row in sheet.iter_rows():
            values = []
            has_data = False
            for cell in row:
                text = "" if cell.value is None else md_escape_cell(str(cell.value))
                if text.strip():
                    has_data = True
                values.append(text)
            if has_data:
                used_rows.append(values)
        if not used_rows:
            out_lines.extend(("_Empty sheet_", ""))
            continue
        col_count = max(len(row) for row in used_rows)
        header = used_rows[0] + [""] * (col_count - len(used_rows[0]))
        out_lines.append("| " + " | ".join(header[:col_count]) + " |")
        out_lines.append("| " + " | ".join(["---"] * col_count) + " |")
        for row in used_rows[1:]:
            padded = row + [""] * (col_count - len(row))
            out_lines.append("| " + " | ".join(padded[:col_count]) + " |")
        out_lines.append("")
    return "\n".join(out_lines).strip() + "\n"


def pdf_to_md(pdf_path: Path) -> str:
    import fitz

    doc = fitz.open(str(pdf_path))
    out_lines: List[str] = []
    for page_index, page in enumerate(doc, start=1):
        text = (page.get_text("text") or "").strip()
        out_lines.extend((f"## Page {page_index}", "", text or "_No extractable text on this page_", ""))
    return "\n".join(out_lines).strip() + "\n"


def html_to_md(html_path: Path) -> str:
    try:
        from markitdown import MarkItDown

        result = MarkItDown().convert(str(html_path))
        return (result.text_content or "").strip() + "\n"
    except Exception:
        raw = read_text_file(html_path)
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
        return text + "\n"


def csv_to_md(csv_path: Path) -> str:
    rows = []
    for line in read_text_file(csv_path).splitlines():
        if line.strip():
            rows.append([md_escape_cell(part) for part in line.split(",")])
    if not rows:
        return ""
    col_count = max(len(row) for row in rows)
    header = rows[0] + [""] * (col_count - len(rows[0]))
    out_lines = [
        "| " + " | ".join(header[:col_count]) + " |",
        "| " + " | ".join(["---"] * col_count) + " |",
    ]
    for row in rows[1:]:
        padded = row + [""] * (col_count - len(row))
        out_lines.append("| " + " | ".join(padded[:col_count]) + " |")
    return "\n".join(out_lines).strip() + "\n"


def convert_local(path: Path) -> Tuple[bool, str]:
    ext = path.suffix.lower()
    if ext == ".docx":
        return True, docx_to_md(path)
    if ext == ".pptx":
        return True, pptx_to_md(path)
    if ext == ".xlsx":
        return True, xlsx_to_md(path)
    if ext == ".pdf":
        return True, pdf_to_md(path)
    if ext == ".html":
        return True, html_to_md(path)
    if ext == ".txt":
        return True, read_text_file(path).strip() + "\n"
    if ext == ".csv":
        return True, csv_to_md(path)
    return False, f"[WARN] Unsupported local format: {ext}"


def get_api_key(path: Path) -> str | None:
    path_name = path.name.lower()
    if "openai" in path_name:
        env_key = os.getenv("AUDION_OPENAI_API_KEY", "").strip()
        env_key_file = os.getenv("AUDION_OPENAI_API_KEY_FILE", "").strip()
        fallback_paths = [LEGACY_CONFIG_API_KEY_OPENAI_FILE, API_KEY_OPENAI_FILE]
    elif "gemini" in path_name:
        env_key = os.getenv("AUDION_GEMINI_API_KEY", "").strip()
        env_key_file = os.getenv("AUDION_GEMINI_API_KEY_FILE", "").strip()
        fallback_paths = [LEGACY_CONFIG_API_KEY_GEMINI_FILE, API_KEY_GEMINI_FILE]
    else:
        env_key = ""
        env_key_file = ""
        fallback_paths = []

    if env_key:
        return env_key
    if env_key_file:
        override_path = Path(env_key_file)
        if override_path.exists():
            key = override_path.read_text(encoding="utf-8", errors="ignore").strip()
            if key:
                return key

    if path.exists():
        key = path.read_text(encoding="utf-8", errors="ignore").strip()
        if key:
            return key
    for fallback_path in fallback_paths:
        if fallback_path == path or not fallback_path.exists():
            continue
        key = fallback_path.read_text(encoding="utf-8", errors="ignore").strip()
        if key:
            return key
    return None


def parse_simple_yaml_settings(path: Path) -> dict:
    if not path.exists():
        return {}

    root: dict = {}
    stack: list[tuple[int, dict]] = [(-1, root)]
    for raw_line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if not raw_line.strip() or raw_line.lstrip().startswith("#"):
            continue
        indent = len(raw_line) - len(raw_line.lstrip(" "))
        line = raw_line.strip()
        if ":" not in line:
            continue
        key, value = line.split(":", 1)
        key = key.strip()
        value = strip_yaml_inline_comment(value.strip())

        while stack and indent <= stack[-1][0]:
            stack.pop()
        parent = stack[-1][1]

        if value == "":
            node: dict = {}
            parent[key] = node
            stack.append((indent, node))
            continue

        if value.lower() == "true":
            parent[key] = True
        elif value.lower() == "false":
            parent[key] = False
        else:
            parent[key] = value.strip("\"'")
    return root


def strip_yaml_inline_comment(value: str) -> str:
    if not value or "#" not in value:
        return value

    in_single = False
    in_double = False
    for index, ch in enumerate(value):
        if ch == "'" and not in_double:
            in_single = not in_single
            continue
        if ch == '"' and not in_single:
            in_double = not in_double
            continue
        if ch == "#" and not in_single and not in_double:
            if index == 0 or value[index - 1].isspace():
                return value[:index].rstrip()
    return value


def load_provider_settings() -> dict:
    settings = parse_simple_yaml_settings(LLM_SETTINGS_FILE)
    providers = settings.get("providers", {}) if isinstance(settings, dict) else {}
    if not isinstance(providers, dict):
        providers = {}
    return providers


def load_llm_settings() -> dict:
    settings = parse_simple_yaml_settings(LLM_SETTINGS_FILE)
    return settings if isinstance(settings, dict) else {}


def _to_int(value, default: int) -> int:
    try:
        return int(str(value).strip())
    except Exception:
        return default


def _to_str(value, default: str) -> str:
    text = str(value).strip() if value is not None else ""
    return text or default


def _to_float(value, default: float) -> float:
    try:
        return float(str(value).strip())
    except Exception:
        return default


def _to_bool(value, default: bool) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "on"}:
        return True
    if text in {"0", "false", "no", "off"}:
        return False
    return default


def load_ocr_settings() -> dict:
    settings = load_llm_settings().get("ocr", {})
    if not isinstance(settings, dict):
        settings = {}

    active_profile = _to_str(settings.get("active_profile"), "")
    profiles = settings.get("profiles", {})
    if isinstance(profiles, dict) and active_profile and isinstance(profiles.get(active_profile), dict):
        merged_settings = dict(profiles.get(active_profile, {}))
    else:
        merged_settings = {k: v for k, v in settings.items() if k not in {"profiles", "active_profile"}}

    openai_cfg = merged_settings.get("openai", {}) if isinstance(merged_settings.get("openai", {}), dict) else {}
    gemini_cfg = merged_settings.get("gemini", {}) if isinstance(merged_settings.get("gemini", {}), dict) else {}
    photo_cfg = merged_settings.get("photo_override", {}) if isinstance(merged_settings.get("photo_override", {}), dict) else {}
    return {
        "active_profile": active_profile or "inline",
        "language_hint": _to_str(merged_settings.get("language_hint"), "ru,en,de"),
        "prompt_file": _to_str(merged_settings.get("prompt_file"), "ocr_prompt_primary_ru.md"),
        "entity_hints_enabled": _to_bool(merged_settings.get("entity_hints_enabled"), True),
        "entity_hints_file": _to_str(merged_settings.get("entity_hints_file"), "ocr_named_entities_ru.yaml"),
        "entity_hints_max_chars": max(0, _to_int(merged_settings.get("entity_hints_max_chars"), 4000)),
        "second_pass_enabled": _to_bool(merged_settings.get("second_pass_enabled"), True),
        "second_pass_if_cyrillic": _to_bool(merged_settings.get("second_pass_if_cyrillic"), True),
        "second_pass_prompt_file": _to_str(merged_settings.get("second_pass_prompt_file"), "ocr_prompt_refine_ru.md"),
        "second_pass_min_cyrillic_ratio": max(0.0, min(1.0, _to_float(merged_settings.get("second_pass_min_cyrillic_ratio"), 0.18))),
        "numeric_mode_threshold": max(0.0, min(1.0, _to_float(merged_settings.get("numeric_mode_threshold"), 0.95))),
        "pdf_render_dpi": _to_int(merged_settings.get("pdf_render_dpi"), 300),
        "pages_per_batch": max(1, _to_int(merged_settings.get("pages_per_batch"), 4)),
        "openai_pages_per_batch": max(1, _to_int(openai_cfg.get("pages_per_batch"), _to_int(merged_settings.get("pages_per_batch"), 4))),
        "openai_fallback_to_gemini_on_region_error": _to_bool(openai_cfg.get("fallback_to_gemini_on_region_error"), True),
        "gemini_pages_per_batch": max(1, _to_int(gemini_cfg.get("pages_per_batch"), _to_int(merged_settings.get("pages_per_batch"), 4))),
        "gemini_fallback_to_openai_on_region_error": _to_bool(gemini_cfg.get("fallback_to_openai_on_region_error"), True),
        "image_min_long_side": max(256, _to_int(merged_settings.get("image_min_long_side"), 1600)),
        "image_max_long_side": max(512, _to_int(merged_settings.get("image_max_long_side"), 2800)),
        "image_format": _to_str(merged_settings.get("image_format"), "PNG").upper(),
        "export_prepared_images": _to_bool(merged_settings.get("export_prepared_images"), False),
        "auto_photo_mode": _to_bool(merged_settings.get("auto_photo_mode"), True),
        "autocontrast": _to_bool(merged_settings.get("autocontrast"), True),
        "grayscale": _to_bool(merged_settings.get("grayscale"), False),
        "contrast": max(0.5, _to_float(merged_settings.get("contrast"), 1.0)),
        "sharpness": max(0.5, _to_float(merged_settings.get("sharpness"), 1.0)),
        "median_filter_size": max(0, _to_int(merged_settings.get("median_filter_size"), 0)),
        "photo_autocontrast": _to_bool(photo_cfg.get("autocontrast"), True),
        "photo_grayscale": _to_bool(photo_cfg.get("grayscale"), True),
        "photo_contrast": max(0.5, _to_float(photo_cfg.get("contrast"), 1.08)),
        "photo_sharpness": max(0.5, _to_float(photo_cfg.get("sharpness"), 1.06)),
        "photo_median_filter_size": max(0, _to_int(photo_cfg.get("median_filter_size"), 3)),
        "openai_timeout_sec": _to_int(openai_cfg.get("timeout_sec"), 180),
        "openai_max_retries": max(1, _to_int(openai_cfg.get("max_retries"), 4)),
        "openai_reasoning_effort": _to_str(openai_cfg.get("reasoning_effort"), "minimal"),
        "openai_verbosity": _to_str(openai_cfg.get("verbosity"), "medium"),
        "gemini_timeout_sec": _to_int(gemini_cfg.get("timeout_sec"), 30),
        "gemini_file_timeout_sec": _to_int(gemini_cfg.get("file_timeout_sec"), _to_int(gemini_cfg.get("timeout_sec"), 30)),
        "gemini_max_retries": max(1, _to_int(gemini_cfg.get("max_retries"), 4)),
    }


def load_primary_ocr_prompt(ocr_settings: dict) -> str:
    path = resolve_config_file(_to_str(ocr_settings.get("prompt_file"), "ocr_prompt_primary_ru.md"), "ocr_prompt_primary_ru.md")
    return ensure_prompt_template(path, DEFAULT_OCR_PROMPT)


def load_second_pass_prompt(ocr_settings: dict) -> str:
    path = resolve_config_file(_to_str(ocr_settings.get("second_pass_prompt_file"), "ocr_prompt_refine_ru.md"), "ocr_prompt_refine_ru.md")
    return ensure_prompt_template(path, DEFAULT_OCR_REFINE_PROMPT)


def load_entity_hints(ocr_settings: dict) -> str:
    if not _to_bool(ocr_settings.get("entity_hints_enabled"), True):
        return "none"
    path = resolve_config_file(_to_str(ocr_settings.get("entity_hints_file"), "ocr_named_entities_ru.yaml"), "ocr_named_entities_ru.yaml")
    text = ensure_prompt_template(path, DEFAULT_ENTITY_HINTS)
    if path.suffix.lower() in {".yaml", ".yml"}:
        try:
            parsed = parse_entity_hints_yaml(text)
            text = format_entity_hints_text(parsed)
        except Exception as exc:
            text = f"[WARN] Failed to parse entity glossary YAML: {exc}\n\n{text}"
    max_chars = max(0, _to_int(ocr_settings.get("entity_hints_max_chars"), 4000))
    if max_chars and len(text) > max_chars:
        text = text[:max_chars].rstrip() + "\n..."
    return text.strip() or "none"


def build_primary_ocr_prompt(source_file: Path, *, chunk_index: int, chunk_count: int, ocr_settings: dict) -> str:
    template = load_primary_ocr_prompt(ocr_settings)
    entity_hints = load_entity_hints(ocr_settings)
    values = {
        "language_hint": _to_str(ocr_settings.get("language_hint"), "ru,en,de"),
        "document_name": source_file.name,
        "chunk_index": str(chunk_index),
        "chunk_count": str(chunk_count),
        "numeric_mode_hint": "auto-detect",
        "entity_hints": entity_hints,
    }
    prompt = render_prompt_template(template, values)
    return (
        f"{prompt}\n\n"
        f"Документ: {source_file.name}\n"
        f"Батч страниц: {chunk_index} из {chunk_count}.\n"
        f"Языки документа: {values['language_hint']}.\n"
        "Не добавляй лишние символы и не галлюцинируй знаки, которых нет на изображении.\n"
        f"Подсказки по именованным сущностям:\n{entity_hints}\n"
        "Сохраняй непрерывность текста между страницами."
    )


def build_second_pass_prompt(
    source_file: Path,
    *,
    chunk_index: int,
    chunk_count: int,
    draft_markdown: str,
    numeric_mode: bool,
    ocr_settings: dict,
) -> str:
    template = load_second_pass_prompt(ocr_settings)
    entity_hints = load_entity_hints(ocr_settings)
    values = {
        "language_hint": _to_str(ocr_settings.get("language_hint"), "ru,en,de"),
        "document_name": source_file.name,
        "chunk_index": str(chunk_index),
        "chunk_count": str(chunk_count),
        "draft_markdown": draft_markdown.strip() or "_empty_",
        "numeric_mode_hint": "yes" if numeric_mode else "no",
        "entity_hints": entity_hints,
    }
    return render_prompt_template(template, values)


def should_run_second_pass(source_file: Path, first_pass_text: str, ocr_settings: dict) -> bool:
    if not _to_bool(ocr_settings.get("second_pass_enabled"), True):
        return False
    if not _to_bool(ocr_settings.get("second_pass_if_cyrillic"), True):
        return True
    threshold = _to_float(ocr_settings.get("second_pass_min_cyrillic_ratio"), 0.18)
    if contains_cyrillic(source_file.name):
        return True
    return cyrillic_ratio(first_pass_text) >= threshold


def is_numeric_heavy_text(text: str, ocr_settings: dict) -> bool:
    threshold = _to_float(ocr_settings.get("numeric_mode_threshold"), 0.95)
    return numeric_signal_ratio(text) >= threshold


def is_openai_region_error(exc: Exception) -> bool:
    msg = str(exc).lower()
    return (
        "unsupported_country_region_territory" in msg
        or "country, region, or territory not supported" in msg
    )


def is_gemini_region_error(exc: Exception) -> bool:
    msg = str(exc).lower()
    return (
        "unsupported_country_region_territory" in msg
        or "country, region, or territory not supported" in msg
        or "user location is not supported" in msg
        or "location is not supported" in msg
        or ("403" in msg and "permission denied" in msg)
        or ("403" in msg and "forbidden" in msg)
    )


def init_openai_client():
    from openai import OpenAI

    api_key = get_api_key(CONFIG_API_KEY_OPENAI_FILE)
    if not api_key:
        raise RuntimeError(f"OpenAI API key not found in: {CONFIG_API_KEY_OPENAI_FILE}")
    return OpenAI(api_key=api_key)


def init_gemini_client():
    from google import genai

    api_key = get_api_key(CONFIG_API_KEY_GEMINI_FILE)
    if not api_key:
        raise RuntimeError(f"Gemini API key not found in: {CONFIG_API_KEY_GEMINI_FILE}")
    cached_client = getattr(_GEMINI_CLIENTS, "client", None)
    cached_key = getattr(_GEMINI_CLIENTS, "api_key", None)
    if cached_client is not None and cached_key == api_key:
        return cached_client
    client = genai.Client(api_key=api_key)
    _GEMINI_CLIENTS.client = client
    _GEMINI_CLIENTS.api_key = api_key
    return client


def file_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def render_pdf_pages(pdf_path: Path, temp_dir: Path, dpi: int = 200) -> list[str]:
    import fitz

    doc = fitz.open(str(pdf_path))
    paths: list[str] = []
    scale = dpi / 72.0
    matrix = fitz.Matrix(scale, scale)
    for idx, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        out_path = temp_dir / f"{pdf_path.stem}_page_{idx:04d}.png"
        pix.save(str(out_path))
        paths.append(str(out_path))
    return paths


def looks_like_camera_photo(source_path: Path, img: Image.Image) -> bool:
    try:
        exif = img.getexif()
    except Exception:
        exif = None
    if not exif:
        return False
    if any(tag in exif for tag in CAMERA_EXIF_TAGS):
        return True
    return source_path.suffix.lower() in {".jpg", ".jpeg", ".tif", ".tiff"} and len(exif) >= 6


def normalize_image_for_ocr(
    source_path: Path,
    temp_dir: Path,
    *,
    image_format: str,
    min_long_side: int,
    max_long_side: int,
    autocontrast: bool,
    grayscale: bool,
    contrast: float,
    sharpness: float,
    median_filter_size: int,
    auto_photo_mode: bool,
    photo_autocontrast: bool,
    photo_grayscale: bool,
    photo_contrast: float,
    photo_sharpness: float,
    photo_median_filter_size: int,
) -> Path:
    with Image.open(source_path) as img:
        img = ImageOps.exif_transpose(img)
        if auto_photo_mode and looks_like_camera_photo(source_path, img):
            autocontrast = photo_autocontrast
            grayscale = photo_grayscale
            contrast = photo_contrast
            sharpness = photo_sharpness
            median_filter_size = photo_median_filter_size
        if grayscale:
            img = img.convert("L")
        elif img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        long_side = max(img.size)
        target_long_side = long_side
        if long_side < min_long_side:
            target_long_side = min_long_side
        elif long_side > max_long_side:
            target_long_side = max_long_side

        if target_long_side != long_side:
            scale = target_long_side / float(long_side)
            new_size = (
                max(1, int(round(img.size[0] * scale))),
                max(1, int(round(img.size[1] * scale))),
            )
            img = img.resize(new_size, Image.Resampling.LANCZOS)

        if autocontrast:
            img = ImageOps.autocontrast(img)
        if contrast != 1.0:
            img = ImageEnhance.Contrast(img).enhance(contrast)
        if sharpness != 1.0:
            img = ImageEnhance.Sharpness(img).enhance(sharpness)
        if median_filter_size >= 3:
            img = img.filter(ImageFilter.MedianFilter(size=median_filter_size))

        suffix = ".png" if image_format == "PNG" else ".jpg"
        out_path = temp_dir / f"{source_path.stem}_ocr{suffix}"
        save_kwargs = {"format": image_format}
        if image_format == "JPEG":
            if img.mode != "RGB":
                img = img.convert("RGB")
            save_kwargs.update({"quality": 92, "optimize": True})
        img.save(out_path, **save_kwargs)
        return out_path


def make_image_inputs(source_file: Path, temp_dir: Path, ocr_settings: dict, ocr_prep: str) -> list[str]:
    image_format = _to_str(ocr_settings.get("image_format"), "PNG").upper()
    min_long_side = _to_int(ocr_settings.get("image_min_long_side"), 1600)
    max_long_side = _to_int(ocr_settings.get("image_max_long_side"), 2800)
    autocontrast = _to_bool(ocr_settings.get("autocontrast"), True)
    grayscale = _to_bool(ocr_settings.get("grayscale"), False)
    contrast = _to_float(ocr_settings.get("contrast"), 1.0)
    sharpness = _to_float(ocr_settings.get("sharpness"), 1.0)
    median_filter_size = _to_int(ocr_settings.get("median_filter_size"), 0)
    auto_photo_mode = _to_bool(ocr_settings.get("auto_photo_mode"), True)
    if ocr_prep == "native":
        auto_photo_mode = False
    elif ocr_prep == "exif-auto":
        auto_photo_mode = _to_bool(ocr_settings.get("auto_photo_mode"), True)
    else:
        raise RuntimeError(f"Unsupported OCR prep mode: {ocr_prep}")
    photo_autocontrast = _to_bool(ocr_settings.get("photo_autocontrast"), True)
    photo_grayscale = _to_bool(ocr_settings.get("photo_grayscale"), True)
    photo_contrast = _to_float(ocr_settings.get("photo_contrast"), 1.08)
    photo_sharpness = _to_float(ocr_settings.get("photo_sharpness"), 1.06)
    photo_median_filter_size = _to_int(ocr_settings.get("photo_median_filter_size"), 3)
    if source_file.suffix.lower() == ".pdf":
        rendered = render_pdf_pages(source_file, temp_dir, dpi=_to_int(ocr_settings.get("pdf_render_dpi"), 300))
        return [
            str(
                normalize_image_for_ocr(
                    Path(rendered_path),
                    temp_dir,
                    image_format=image_format,
                    min_long_side=min_long_side,
                    max_long_side=max_long_side,
                    autocontrast=autocontrast,
                    grayscale=grayscale,
                    contrast=contrast,
                    sharpness=sharpness,
                    median_filter_size=median_filter_size,
                    auto_photo_mode=False,
                    photo_autocontrast=photo_autocontrast,
                    photo_grayscale=photo_grayscale,
                    photo_contrast=photo_contrast,
                    photo_sharpness=photo_sharpness,
                    photo_median_filter_size=photo_median_filter_size,
                )
            )
            for rendered_path in rendered
        ]
    normalized = normalize_image_for_ocr(
        source_file,
        temp_dir,
        image_format=image_format,
        min_long_side=min_long_side,
        max_long_side=max_long_side,
        autocontrast=autocontrast,
        grayscale=grayscale,
        contrast=contrast,
        sharpness=sharpness,
        median_filter_size=median_filter_size,
        auto_photo_mode=auto_photo_mode,
        photo_autocontrast=photo_autocontrast,
        photo_grayscale=photo_grayscale,
        photo_contrast=photo_contrast,
        photo_sharpness=photo_sharpness,
        photo_median_filter_size=photo_median_filter_size,
    )
    return [str(normalized)]


def export_prepared_images(source_file: Path, prepared_image_paths: list[str]) -> None:
    base_target = BASE_DIR / "output" / "_ocr_debug" / normalized_source_relative_path(source_file).with_suffix("")
    for index, prepared_path in enumerate(prepared_image_paths, start=1):
        prepared = Path(prepared_path)
        if len(prepared_image_paths) == 1:
            target = base_target.with_name(base_target.name + "_ocr").with_suffix(prepared.suffix)
        else:
            target = base_target / f"page_{index:04d}{prepared.suffix}"
        ensure_parent_dir(target)
        shutil.copy2(prepared, target)


def chunk_list(items: list[str], chunk_size: int) -> list[list[str]]:
    return [items[i:i + chunk_size] for i in range(0, len(items), chunk_size)]


def normalize_markdown(text: str) -> str:
    value = (text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def cleanup_temp_dir(path: Path, *, retries: int = 5, delay_sec: float = 0.35) -> None:
    for attempt in range(1, retries + 1):
        try:
            gc.collect()
            shutil.rmtree(path)
            return
        except FileNotFoundError:
            return
        except PermissionError:
            if attempt >= retries:
                print(f"[WARN] Temporary OCR folder was left for manual cleanup: {path}")
                return
            time.sleep(delay_sec * attempt)


def resolve_model(provider: str, model: str) -> str:
    if model and model.lower() != "auto":
        return normalize_model_name(provider, model)
    provider_settings = load_provider_settings().get(provider, {})
    if isinstance(provider_settings, dict):
        configured_model = str(provider_settings.get("model", "")).strip()
        if configured_model:
            return normalize_model_name(provider, configured_model)
    if provider == "openai":
        return "gpt-5.4"
    if provider == "gemini":
        return DEFAULT_GEMINI_MODEL_CHAIN[0]
    raise RuntimeError(f"Unsupported provider for model resolution: {provider}")


def gemini_model_failover_chain(model: str, resolved_model: str) -> list[str] | None:
    model_text = str(model or "").strip().lower()
    if model_text and model_text != "auto":
        return None
    normalized = normalize_model_name("gemini", resolved_model)
    chain = [normalize_model_name("gemini", item) for item in DEFAULT_GEMINI_MODEL_CHAIN]
    if normalized in chain:
        return chain[chain.index(normalized):]
    return [normalized, *chain]


def normalize_model_name(provider: str, model: str) -> str:
    text = str(model or "").strip()
    if not text:
        return text
    if provider == "openai":
        return text.lower()
    if provider == "gemini":
        text = text.removeprefix("models/")
        return text.lower()
    return text


def ensure_provider_enabled(provider: str) -> None:
    provider_settings = load_provider_settings().get(provider, {})
    if isinstance(provider_settings, dict):
        enabled = provider_settings.get("enabled", True)
        if enabled is False:
            raise RuntimeError(f"Provider '{provider}' is disabled in {LLM_SETTINGS_FILE}")


def sort_files_for_processing(files: list[Path]) -> list[Path]:
    return sorted(files, key=lambda path: str(normalized_source_relative_path(path)).lower())


def merged_output_path(provider: str) -> Path:
    return BASE_DIR / "output" / f"images_merged_{provider}.md"


def call_openai_vision_chunk(
    client,
    *,
    model: str,
    source_file: Path,
    chunk: list[str],
    prompt: str,
    doc_hash: str,
    chunk_index: int,
    ocr_settings: dict,
    pass_label: str,
) -> str:
    instructions = "Convert document images into clean Markdown."
    if pass_label == "refine":
        instructions = "Correct OCR draft using the same document images and return only corrected Markdown."
    text, _usage, _tier = openai_call_markdown_vision(
        client,
        model=model,
        instructions=instructions,
        user_prompt=prompt,
        image_paths=chunk,
        reasoning_effort=_to_str(ocr_settings.get("openai_reasoning_effort"), "minimal"),
        max_output_tokens=16000,
        timeout_sec=float(_to_int(ocr_settings.get("openai_timeout_sec"), 180)),
        max_retries=_to_int(ocr_settings.get("openai_max_retries"), 4),
        service_tier="auto",
        use_idempotency=True,
        doc_hash=doc_hash,
        chunk_index=chunk_index,
        verbosity=_to_str(ocr_settings.get("openai_verbosity"), "medium"),
    )
    return normalize_markdown(text)


def call_gemini_vision_chunk(
    client,
    *,
    model: str,
    model_chain: list[str] | None = None,
    chunk: list[str],
    prompt: str,
    ocr_settings: dict,
    deadline_monotonic: float,
    pass_label: str,
) -> str:
    per_request_timeout_sec = max(1, _to_int(ocr_settings.get("gemini_timeout_sec"), 30))
    remaining_sec = deadline_monotonic - time.monotonic()
    if remaining_sec <= 0:
        raise RuntimeError("Gemini OCR timed out before the next request could start.")
    system_instruction = "Convert document images into clean Markdown."
    if pass_label == "refine":
        system_instruction = "Correct OCR draft using the same document images and return only corrected Markdown."
    text, _usage, _tier = gemini_call_markdown_vision(
        client,
        model=model,
        model_chain=model_chain,
        system_instruction=system_instruction,
        user_prompt=prompt,
        image_paths=chunk,
        temperature=0.0,
        timeout_sec=min(per_request_timeout_sec, max(1, int(remaining_sec))),
        deadline_monotonic=deadline_monotonic,
        max_retries=_to_int(ocr_settings.get("gemini_max_retries"), 4),
        sleep_after_sec=0.0,
    )
    return normalize_markdown(text)


def call_openai_ocr(source_file: Path, image_paths: list[str], model: str) -> str:
    ocr_settings = load_ocr_settings()
    client = init_openai_client()
    doc_hash = file_sha256(source_file)
    chunks = chunk_list(image_paths, _to_int(ocr_settings.get("openai_pages_per_batch"), _to_int(ocr_settings.get("pages_per_batch"), 4)))
    outputs: list[str] = []

    for idx, chunk in enumerate(chunks, start=1):
        primary_prompt = build_primary_ocr_prompt(source_file, chunk_index=idx, chunk_count=len(chunks), ocr_settings=ocr_settings)
        first_pass = call_openai_vision_chunk(
            client,
            model=model,
            source_file=source_file,
            chunk=chunk,
            prompt=primary_prompt,
            doc_hash=doc_hash,
            chunk_index=(idx - 1) * 2,
            ocr_settings=ocr_settings,
            pass_label="primary",
        )
        final_pass = first_pass
        if should_run_second_pass(source_file, first_pass, ocr_settings):
            numeric_mode = is_numeric_heavy_text(first_pass, ocr_settings)
            refine_prompt = build_second_pass_prompt(
                source_file,
                chunk_index=idx,
                chunk_count=len(chunks),
                draft_markdown=first_pass,
                numeric_mode=numeric_mode,
                ocr_settings=ocr_settings,
            )
            final_pass = call_openai_vision_chunk(
                client,
                model=model,
                source_file=source_file,
                chunk=chunk,
                prompt=refine_prompt,
                doc_hash=doc_hash,
                chunk_index=(idx - 1) * 2 + 1,
                ocr_settings=ocr_settings,
                pass_label="refine",
            )
        outputs.append(final_pass)

    return "\n\n".join(part for part in outputs if part).strip() + "\n"


def call_gemini_ocr(source_file: Path, image_paths: list[str], model: str, model_chain: list[str] | None = None) -> str:
    ocr_settings = load_ocr_settings()
    client = init_gemini_client()
    chunks = chunk_list(image_paths, _to_int(ocr_settings.get("gemini_pages_per_batch"), _to_int(ocr_settings.get("pages_per_batch"), 4)))
    outputs: list[str] = []
    file_timeout_sec = max(1, _to_int(ocr_settings.get("gemini_file_timeout_sec"), 30))
    deadline_monotonic = time.monotonic() + file_timeout_sec

    for idx, chunk in enumerate(chunks, start=1):
        if deadline_monotonic - time.monotonic() <= 0:
            raise RuntimeError(f"Gemini OCR timed out after {file_timeout_sec}s for file: {source_file.name}")
        primary_prompt = build_primary_ocr_prompt(source_file, chunk_index=idx, chunk_count=len(chunks), ocr_settings=ocr_settings)
        first_pass = call_gemini_vision_chunk(
            client,
            model=model,
            model_chain=model_chain,
            chunk=chunk,
            prompt=primary_prompt,
            ocr_settings=ocr_settings,
            deadline_monotonic=deadline_monotonic,
            pass_label="primary",
        )
        final_pass = first_pass
        if should_run_second_pass(source_file, first_pass, ocr_settings):
            numeric_mode = is_numeric_heavy_text(first_pass, ocr_settings)
            refine_prompt = build_second_pass_prompt(
                source_file,
                chunk_index=idx,
                chunk_count=len(chunks),
                draft_markdown=first_pass,
                numeric_mode=numeric_mode,
                ocr_settings=ocr_settings,
            )
            final_pass = call_gemini_vision_chunk(
                client,
                model=model,
                model_chain=model_chain,
                chunk=chunk,
                prompt=refine_prompt,
                ocr_settings=ocr_settings,
                deadline_monotonic=deadline_monotonic,
                pass_label="refine",
            )
        outputs.append(final_pass)

    return "\n\n".join(part for part in outputs if part).strip() + "\n"


def ai_ocr_text(source_file: Path, provider: str, model: str, ocr_prep: str) -> tuple[str, str]:
    temp_dir = Path(tempfile.mkdtemp(prefix="audion_ocr_", dir=str(TEMP_DIR)))
    try:
        ocr_settings = load_ocr_settings()
        image_paths = make_image_inputs(source_file, temp_dir, ocr_settings, ocr_prep)
        if not image_paths:
            raise RuntimeError("No rendered pages/images were produced for OCR.")
        if _to_bool(ocr_settings.get("export_prepared_images"), False):
            export_prepared_images(source_file, image_paths)
        attempted_providers: set[str] = set()
        current_provider = provider
        current_model = model
        fallback_from: str | None = None

        while True:
            ensure_provider_enabled(current_provider)
            attempted_providers.add(current_provider)

            try:
                if current_provider == "openai":
                    text = call_openai_ocr(source_file, image_paths, resolve_model("openai", current_model))
                elif current_provider == "gemini":
                    resolved_model = resolve_model("gemini", current_model)
                    text = call_gemini_ocr(
                        source_file,
                        image_paths,
                        resolved_model,
                        model_chain=gemini_model_failover_chain(current_model, resolved_model),
                    )
                else:
                    raise RuntimeError(f"Unsupported AI OCR provider: {current_provider}")

                if fallback_from:
                    return f"{current_provider} (fallback from {fallback_from})", text
                return current_provider, text

            except Exception as exc:
                next_provider: str | None = None

                if (
                    current_provider == "openai"
                    and _to_bool(ocr_settings.get("openai_fallback_to_gemini_on_region_error"), True)
                    and is_openai_region_error(exc)
                    and "gemini" not in attempted_providers
                ):
                    next_provider = "gemini"
                elif (
                    current_provider == "gemini"
                    and _to_bool(ocr_settings.get("gemini_fallback_to_openai_on_region_error"), True)
                    and is_gemini_region_error(exc)
                    and "openai" not in attempted_providers
                ):
                    next_provider = "openai"

                if next_provider is None:
                    raise

                print(
                    f"[FALLBACK] {current_provider.capitalize()} region-restricted for "
                    f"{source_file.name}; retrying via {next_provider.capitalize()} "
                    f"({resolve_model(next_provider, 'auto')})."
                )
                previous_provider = current_provider
                previous_error = exc
                fallback_from = current_provider
                current_provider = next_provider
                current_model = "auto"

                try:
                    ensure_provider_enabled(current_provider)
                    attempted_providers.add(current_provider)
                    if current_provider == "openai":
                        text = call_openai_ocr(source_file, image_paths, resolve_model("openai", current_model))
                    elif current_provider == "gemini":
                        resolved_model = resolve_model("gemini", current_model)
                        text = call_gemini_ocr(
                            source_file,
                            image_paths,
                            resolved_model,
                            model_chain=gemini_model_failover_chain(current_model, resolved_model),
                        )
                    else:
                        raise RuntimeError(f"Unsupported AI OCR provider: {current_provider}")
                    return f"{current_provider} (fallback from {fallback_from})", text
                except Exception as fallback_exc:
                    raise RuntimeError(
                        f"{previous_provider.capitalize()} region-restricted for {source_file.name}; "
                        f"fallback via {current_provider.capitalize()} also failed: {fallback_exc}"
                    ) from fallback_exc
    finally:
        cleanup_temp_dir(temp_dir)


def ai_ocr_to_md(source_file: Path, out_file: Path, provider: str, model: str, ocr_prep: str) -> tuple[str, str]:
    used_provider, text = ai_ocr_text(source_file, provider=provider, model=model, ocr_prep=ocr_prep)
    ensure_parent_dir(out_file)
    out_file.write_text(text, encoding="utf-8")
    return used_provider, text


def process_documents(force_local: bool, model: str, provider: str, mode: str, ocr_prep: str) -> None:
    ensure_project_dirs()
    supported_exts = LOCAL_EXTS.union(AI_EXTS)
    all_files = iter_project_files(supported_exts)
    all_files = sort_files_for_processing(all_files)

    if not all_files:
        print(f"[INFO] No supported files found in project root: {BASE_DIR}")
        return

    if force_local:
        all_files = [file for file in all_files if file.suffix.lower() in LOCAL_EXTS]

    if mode == "images_per_file":
        all_files = [file for file in all_files if file.suffix.lower() in IMAGE_EXTS]
    elif mode == "images_merged":
        all_files = [file for file in all_files if file.suffix.lower() in IMAGE_EXTS]

    if not all_files:
        print(f"[INFO] No supported files found for mode '{mode}' in project root: {BASE_DIR}")
        return

    print("=" * 60)
    if force_local:
        print("AUDION OFFICE OCR AI : FAST LOCAL EXTRACTOR")
    elif mode == "images_per_file":
        print("AUDION OFFICE OCR AI : IMAGES TO MD")
    elif mode == "images_merged":
        print("AUDION OFFICE OCR AI : IMAGES TO ONE MD")
    else:
        print("AUDION OFFICE OCR AI : AI OCR EXTRACTOR")
    print("=" * 60)
    print(f"Found {len(all_files)} files to process.\n")
    if not force_local:
        print(f"OCR prep mode: {ocr_prep}\n")
    failure_count = 0

    if mode == "images_merged":
        merged_parts: list[str] = []
        for source_file in tqdm(all_files, desc="Extracting", unit="image"):
            rel_path = normalized_source_relative_path(source_file)
            tqdm.write(f"\n[AI OCR MERGED] {source_file.name}")
            try:
                used_provider, text = ai_ocr_text(source_file, provider=provider, model=model, ocr_prep=ocr_prep)
                normalized = normalize_markdown(text)
                if normalized:
                    merged_parts.append(f"<!-- source: {rel_path.as_posix()} -->\n\n{normalized}")
                tqdm.write(f"   [OK] Added via {used_provider}: {rel_path}")
            except Exception as exc:
                failure_count += 1
                tqdm.write(f"   [ERROR] AI OCR failed for {source_file.name}: {exc}")

        if merged_parts:
            out_file = merged_output_path(provider)
            ensure_parent_dir(out_file)
            out_file.write_text("\n\n---\n\n".join(merged_parts).strip() + "\n", encoding="utf-8")
            print(f"\n[OK] Merged markdown saved: {out_file}")
        else:
            print("\n[WARN] No merged markdown content was produced.")
        if failure_count:
            print(f"\n[WARN] Extraction finished with {failure_count} failed file(s).")
        else:
            print("\n[OK] Extraction complete.")
        return

    for source_file in tqdm(all_files, desc="Extracting", unit="file"):
        ext = source_file.suffix.lower()
        rel_path = normalized_source_relative_path(source_file)
        out_file = output_path_for(source_file, ".md")
        ensure_parent_dir(out_file)

        if force_local or mode == "all" and ext not in AI_EXTS:
            tqdm.write(f"\n[LOCAL] {source_file.name}")
            try:
                ok, md_or_message = convert_local(source_file)
                if ok:
                    out_file.write_text(md_or_message, encoding="utf-8")
                    tqdm.write(f"   [OK] Saved: {out_file}")
                else:
                    failure_count += 1
                    tqdm.write(f"   {md_or_message}")
            except Exception as exc:
                failure_count += 1
                tqdm.write(f"   [ERROR] Local convert failed: {exc}")
            continue

        tqdm.write(f"\n[AI OCR] {source_file.name}")
        try:
            used_provider, _text = ai_ocr_to_md(
                source_file,
                out_file,
                provider=provider,
                model=model,
                ocr_prep=ocr_prep,
            )
            tqdm.write(f"   [OK] Saved via {used_provider}: {out_file}")
        except Exception as exc:
            failure_count += 1
            tqdm.write(f"   [ERROR] AI OCR failed for {source_file.name}: {exc}")

    if failure_count:
        print(f"\n[WARN] Extraction finished with {failure_count} failed file(s).")
    else:
        print("\n[OK] Extraction complete.")


def main() -> None:
    os.system("cls" if os.name == "nt" else "clear")

    parser = argparse.ArgumentParser(description="Extract documents to Markdown")
    parser.add_argument("--local", action="store_true", help="Force local extraction for all files")
    parser.add_argument("--ai", action="store_true", help="Use AI OCR for supported files")
    parser.add_argument("--provider", default="gemini", choices=["openai", "gemini"], help="AI OCR provider")
    parser.add_argument("--model", default="auto", help="Model name for selected provider, or auto")
    parser.add_argument(
        "--mode",
        default="all",
        choices=["all", "images_per_file", "images_merged"],
        help="Extraction mode",
    )
    parser.add_argument(
        "--ocr-prep",
        default="native",
        choices=["native", "exif-auto"],
        help="Image preparation mode for AI OCR",
    )
    args = parser.parse_args()

    force_local = bool(args.local)
    if not args.local and not args.ai:
        force_local = False
    if force_local and args.mode != "all":
        raise SystemExit("Image OCR modes require --ai, not --local.")

    process_documents(
        force_local=force_local,
        model=args.model,
        provider=args.provider,
        mode=args.mode,
        ocr_prep=args.ocr_prep,
    )


if __name__ == "__main__":
    main()
