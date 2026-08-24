from __future__ import annotations

import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from project_paths import ensure_project_dirs, iter_project_files, normalized_source_relative_path, output_path_for


def clean_markdown_text(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    return text.strip()


def apply_font_style(paragraph, font_size: int = 18, bold: bool = False) -> None:
    for run in paragraph.runs:
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = bold


def compile_markdown_to_pptx() -> None:
    ensure_project_dirs()
    md_files = iter_project_files({".md"})
    if not md_files:
        print("[INFO] No .md files found in project root.")
        return

    print("=== AUDION OFFICE OCR AI: PPTX COMPILER ===")

    for md_file in md_files:
        rel_path = normalized_source_relative_path(md_file)
        out_file = output_path_for(md_file, ".pptx")
        out_file.parent.mkdir(parents=True, exist_ok=True)
        print(f"[BUILD] {rel_path} -> PPTX")

        prs = Presentation()
        layout_title = prs.slide_layouts[0]
        layout_content = prs.slide_layouts[1]
        current_slide = None
        text_frame = None

        for raw_line in md_file.read_text(encoding="utf-8").splitlines():
            line = raw_line.strip()
            if not line or line.startswith("|"):
                continue

            if line.startswith("# "):
                current_slide = prs.slides.add_slide(layout_title)
                title_shape = current_slide.shapes.title
                title_shape.text = clean_markdown_text(line[2:])
                apply_font_style(title_shape.text_frame.paragraphs[0], font_size=30, bold=True)
                text_frame = None
                continue

            if line.startswith("## ") or line.startswith("### "):
                prefix_len = 3 if line.startswith("## ") else 4
                current_slide = prs.slides.add_slide(layout_content)
                title_shape = current_slide.shapes.title
                title_shape.text = clean_markdown_text(line[prefix_len:])
                apply_font_style(title_shape.text_frame.paragraphs[0], font_size=26, bold=True)
                text_frame = current_slide.shapes.placeholders[1].text_frame
                text_frame.clear()
                continue

            if current_slide is None:
                current_slide = prs.slides.add_slide(layout_content)
                text_frame = current_slide.shapes.placeholders[1].text_frame
                text_frame.clear()

            paragraph = text_frame.add_paragraph()
            if line.startswith(("- ", "* ")):
                paragraph.text = clean_markdown_text(line[2:])
                paragraph.level = 0
                apply_font_style(paragraph, font_size=18)
            else:
                paragraph.text = clean_markdown_text(line)
                paragraph.level = 0
                apply_font_style(paragraph, font_size=16)

        prs.save(str(out_file))

    print("\n[OK] Compilation complete.")


if __name__ == "__main__":
    os.system("cls" if os.name == "nt" else "clear")
    compile_markdown_to_pptx()
